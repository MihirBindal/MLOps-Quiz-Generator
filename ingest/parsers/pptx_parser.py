from pptx import Presentation
from PIL import Image, ImageEnhance, ImageOps
from io import BytesIO
import pytesseract

def preprocess_image_for_ocr(img: Image.Image) -> Image.Image:
    """Enhances terminal screenshots and low-res images for better Tesseract OCR."""
    # 1. Convert to grayscale
    img = img.convert('L')
    
    # 2. Upscale the image by 2x (Terminal text is usually too small for OCR)
    new_size = (img.width * 2, img.height * 2)
    img = img.resize(new_size, Image.Resampling.LANCZOS)
    
    # 3. Maximize contrast
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)
    
    # 4. Auto-invert (If it's a dark terminal screenshot, make it black text on white background)
    # We check the average pixel value. If it's dark, we invert it.
    import numpy as np
    if np.mean(np.array(img)) < 127:
        img = ImageOps.invert(img)
        
    return img

def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
            if hasattr(shape, "image"):
                try:
                    image_bytes = shape.image.blob
                    img = Image.open(BytesIO(image_bytes))
                    
                    # Apply our new preprocessing pipeline before OCR
                    enhanced_img = preprocess_image_for_ocr(img)
                    full_text += pytesseract.image_to_string(enhanced_img) + "\n"
                except Exception as e:
                    print(f"Failed to process an image in PPTX: {e}")
                    continue
    return full_text